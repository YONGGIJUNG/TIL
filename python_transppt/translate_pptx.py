# -*- coding: utf-8 -*-

from pptx import Presentation  # python-pptx 라이브러리
import papa                    # 위에서 만든 파파고 라이브러리
# import pypapago
from pypapago import Translator

# import xlrd
import xlwt
import os

# from tkinter import *
import tkinter as tk
from tkinter import ttk
import tkinter.filedialog
from tkinter import messagebox
from tkinter.filedialog import askopenfilename

pa_translator = Translator()

# translator = pypapago.Translator()

'''
번역
shape : pptx 객체
src_lang : 번역할 문자열 언어
tgt_lang : 번역을 요청할 언어
'''
#ko Korean
#en English
#ja Japanese
#zh-CN Chinese
#zh-TW Chinese traditional
#es Spanish
#fr French
#vi Vietnamese
#th Thai
#id Indonesia

# pptx의 text를 속성 그대로 살려두고 글자만 변경
def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    if len(paragraph.runs) == 0:
        return
    p = paragraph._p    
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text



def translate(shape, src_lang, tgt_lang):
    for paragraph in shape.text_frame.paragraphs:
        # 문자열 예외 처리 공백..
        if paragraph.text.isspace():
            continue
        if paragraph.text is '':
            continue             

        # 문자열이 번역할 언어인지 확인 : 언어감지 필요할 경우 
#         if papa.check_language(paragraph.text) == src_lang:
        if src_lang == src_lang:
            # tgt_lang 언어로 번역
#             new_text = papa.translate(paragraph.text, src_lang, tgt_lang)
#             print(paragraph.text)
#             print(shape)
            new_text = pa_translator.translate(paragraph.text, src_lang, tgt_lang)
            # 리턴 값이 none이라면 API 에러로 더이상 진행이 안되므로 바로 리턴
            if new_text is None:
                return False
#             문자열 표시

            with open('{0}_dictionary.txt'.format(os.path.splitext(src_file)[0]),'a', encoding='utf-8') as f:
                print('원문(src) :' + paragraph.text + "  ,   번역(target) : " + new_text+'\n')
                f.write('원문(src) :' + paragraph.text + "  ,   번역(target) : " + new_text+'\n')



            # pptx 해당 문자열 변경
            replace_paragraph_text_retaining_initial_formatting(paragraph, new_text)        
    # 문제없이 번역되었다면 True 리턴
    return True

# 테이블 안에 있는 문자열 번역
def table_func(shape, src_lang, tgt_lang):
    table = shape.table
    # 셀단위로 체크
    for cell in table.iter_cells():
        if translate(cell, src_lang, tgt_lang) is False:
            return False
    return True

# 텍스트 상자 문자열 번역
def text_func(shape, src_lang, tgt_lang):
    return translate(shape, src_lang, tgt_lang)

def run(src_file, tgt_file, src_lang, tgt_lang):                
    # pptx 불러오기
    prs = Presentation(src_file)    
    # 슬라이드 단위로 루프
    for slide in prs.slides:
        # 슬라이드에서 각 객체(모양?) 단위로 루프
        for shape in slide.shapes: 
            # 객체가 테이블이면
            if shape.has_table:
                # 테이블 번역 함수 호출
                if table_func(shape, src_lang, tgt_lang) is False:
                    # 번역함수에서 리턴이 False라면 웹 API에 이상으로 지금까지 작업한 내용 저장하고 리턴
                    prs.save(tgt_file)
                    return False
            # 텍스트 상자면
            if shape.has_text_frame:
                # 텍스트 번역 함수 호출
                if text_func(shape, src_lang, tgt_lang) is False:
                    # 번역함수에서 리턴이 False라면 웹 API에 이상으로 지금까지 작업한 내용 저장하고 리턴
                    prs.save(tgt_file)
                    return False
    # 끝까지 문제없이 번역되었다면 저장하고 리턴
    prs.save(tgt_file)
    return True


if __name__ == "__main__":
    
    app = tk.Tk() 
    app.title("중국빅데이터센터 문서번역기")
    app.geometry('200x300')
    
#     app.withdraw()
#     app.update() 
    src_file = tk.filedialog.askopenfilename() # 번역할 pptx파일  
    app.destroy()
    
    src_lang = 'zh-CN' # 번역할 언어        zh-CN = 중국어 
    tgt_lang = 'ko' # 번역을 요청할 언어 ko = 한국어
#     src_file = '원본.pptx' # 번역할 pptx파일
    tgt_file = os.path.splitext(src_file)[0]+'_번역본.pptx' # 번역되어 저장할 pptx 파일
    
    if os.path.exists('{0}_dictionary.txt'.format(os.path.splitext(src_file)[0])):
        os.remove('{0}_dictionary.txt'.format(os.path.splitext(src_file)[0]))
    else:
        pass
    

#     def callbackFunc(event):
#          print("New Element Selected")

    
#     def clickMe():
#            messagebox.showinfo("번역실행", str.get())
#     src_lang = ttk.Combobox()
#     src_lang['values'] = ('zh-CN', 'ko', 'en','zh-TW')
#     src_lang.grid(column = 0 , row = 0)
#     src_lang.current(0)
# #     action=ttk.Button(text="번역실행", command=clickMe)
# #     action.grid(column=0, row=1)
    

    
    
#     app.mainloop()

    # 번역 시작
    run(src_file, tgt_file, src_lang, tgt_lang)